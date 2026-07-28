"""Inspect table and picture on slide 3, plus check picture on slide 4."""
import os
from pptx import Presentation
from pptx.util import Emu

src = os.path.join(os.environ["TEMP"], "uc_template.pptx")
prs = Presentation(src)

print("=== Slide 3 detail ===")
slide = prs.slides[2]
for shape in slide.shapes:
    print(f"\n{shape.shape_type} name={shape.name!r} left={Emu(shape.left).inches:.2f}in top={Emu(shape.top).inches:.2f}in w={Emu(shape.width).inches:.2f}in h={Emu(shape.height).inches:.2f}in")
    if shape.has_table:
        tbl = shape.table
        print(f"  Table: rows={len(tbl.rows)} cols={len(tbl.columns)}")
        for r_idx, row in enumerate(tbl.rows):
            for c_idx, cell in enumerate(row.cells):
                text = cell.text_frame.text.strip()
                if text:
                    print(f"    [{r_idx},{c_idx}] {text!r}")

print("\n=== Slide 4 detail ===")
slide = prs.slides[3]
for shape in slide.shapes:
    print(f"\n{shape.shape_type} name={shape.name!r} left={Emu(shape.left).inches:.2f}in top={Emu(shape.top).inches:.2f}in w={Emu(shape.width).inches:.2f}in h={Emu(shape.height).inches:.2f}in")
