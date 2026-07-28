"""Build the TFC / One Fiscale integration deck by cloning the Unicredit template
and editing text + swapping in the new architecture / capability images.

Slide mapping (mirrors the Unicredit RAG proposal exactly):
  1. Cover                                (Light Background, Full IC on Right)
  2. Contesto e Sfida                     (Title, 1 Column Body Text)  – 5 h/copy blocks
  3. AI Assistant: Obiettivo & Funzionamento (1_Case Study) – table + right image
  4. Soluzione Proposta – Architettura    (Title, 2 Column Body Text) – text + right image
  5. Casi d'Uso Chiave                    (Title, Subtitle, 1 Column Body Text) – 3 head/copy
  6. Approccio Proof of Concept           (Title Only)
  7. Modello Costi                        (Title, Subtitle, 1 Column Body Text)
  8. End page
"""
from __future__ import annotations

import copy
import os
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

ROOT = Path(__file__).parent.parent
REPORTS = ROOT / "reports"
ASSETS = REPORTS / "assets"

SRC_PPTX = REPORTS / "Wolters Kluwer - Unicredit - RAG Proposal.pptx"
OUT_PPTX = REPORTS / "Wolters Kluwer - TFC - One Fiscale Integration.pptx"

ARCH_PNG = ASSETS / "tfc_architecture.png"
CAPS_PNG = ASSETS / "tfc_capabilities.png"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def find_shape(slide, name: str):
    for shp in slide.shapes:
        if shp.name == name:
            return shp
    raise KeyError(f"Shape {name!r} not found on slide")


def set_text_preserving_format(tf, lines):
    """Replace text of a text_frame with `lines` (list[str]) preserving the
    formatting of the first run of the first paragraph. Extra paragraphs
    inherit the first paragraph's alignment/level and the first run's font.
    """
    if isinstance(lines, str):
        lines = [lines]
    if not lines:
        lines = [""]

    # Capture reference formatting from paragraph 0 / run 0 (if any)
    first_para = tf.paragraphs[0]
    ref_font = None
    ref_align = first_para.alignment
    ref_level = first_para.level
    if first_para.runs:
        ref_font = first_para.runs[0].font

    # Wipe existing paragraphs XML except the first
    p_elems = list(tf._txBody.iter(qn("a:p")))
    for p in p_elems[1:]:
        p.getparent().remove(p)

    # Reset first paragraph
    for r in list(first_para.runs):
        r._r.getparent().remove(r._r)
    first_para.text = lines[0]
    if ref_font is not None and first_para.runs:
        _copy_font(ref_font, first_para.runs[0].font)
    first_para.alignment = ref_align
    first_para.level = ref_level

    for line in lines[1:]:
        new_p = tf.add_paragraph()
        new_p.text = line
        new_p.alignment = ref_align
        new_p.level = ref_level
        if ref_font is not None and new_p.runs:
            _copy_font(ref_font, new_p.runs[0].font)


def _copy_font(src, dst):
    if src.size is not None:
        dst.size = src.size
    if src.name is not None:
        dst.name = src.name
    if src.bold is not None:
        dst.bold = src.bold
    if src.italic is not None:
        dst.italic = src.italic
    try:
        if src.color and src.color.type is not None:
            dst.color.rgb = src.color.rgb
    except Exception:
        pass


def qn(tag):
    # Local import to avoid touching global namespace before lxml is loaded
    from pptx.oxml.ns import qn as _qn
    return _qn(tag)


def set_head_copy(slide, head_name: str, copy_name: str, heading: str, body_lines: list[str]):
    head_shape = find_shape(slide, head_name)
    copy_shape = find_shape(slide, copy_name)
    set_text_preserving_format(head_shape.text_frame, heading)
    set_text_preserving_format(copy_shape.text_frame, body_lines)


def replace_picture(slide, picture_name: str, image_path: Path):
    """Delete the given picture shape and insert a new one at the same
    position and size, keeping z-order approximately the same.
    """
    pic = find_shape(slide, picture_name)
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def set_table_cell(table, row: int, col: int, lines):
    tf = table.cell(row, col).text_frame
    set_text_preserving_format(tf, lines if isinstance(lines, list) else [lines])


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------
COVER_TITLE = "Integrazione One Fiscale su TFC"
COVER_SUB = (
    "Estensione del Tax Control Framework con capacità AI di Wolters Kluwer "
    "One Fiscale: sei assistenti dedicati, orchestrati su FAB, che portano "
    "gap analysis, driver assistant, impact analysis, drafting, Q&A "
    "contestuale e briefing periodici direttamente nell’ambiente TCF del cliente."
)

CONTEXT_TITLE = "Contesto e Sfida"

CONTEXT_HEADER = "Il Tax Control Framework deve evolvere per stare al passo con normativa e operatività fiscale"

CONTEXT_BLOCKS = {
    "h1": (
        "Framework fiscale sempre più articolato",
        "Il Tax Control Framework combina wizard di compilazione secondo Ag. Entrate, RCM, "
        "fattispecie interpretative e strategia fiscale: mantenere allineamento e completezza "
        "diventa oneroso senza un supporto intelligente.",
    ),
    "h2": (
        "Rischi non mappati e valutazioni soggettive",
        "La configurazione di driver di probabilità e impatto è oggi manuale e discrezionale, "
        "esponendo il cliente a rischi non intercettati e a valutazioni disomogenee tra periodi "
        "e aziende del gruppo.",
    ),
    "h4": (
        "Novità normative da tradurre in azioni",
        "Ogni circolare, prassi o articolo di One Fiscale può impattare la RCM: serve un "
        "meccanismo automatico che colleghi la novità normativa ai rischi specifici del cliente "
        "invece di richiedere una revisione integrale.",
    ),
    "h5": (
        "Valorizzare la banca dati One Fiscale nel TCF",
        "One Fiscale contiene una base normativa qualificata e continuamente aggiornata: "
        "integrarla nativamente nel TCF rende ogni funzionalità immediatamente contestualizzata "
        "sul cliente, senza duplicare contenuti né esporli fuori dal perimetro WK.",
    ),
    "h6": (
        "Opportunità di trasformazione operativa",
        "L’adozione di sei agenti AI dedicati abilita un salto di qualità: da compilazione "
        "documentale a supporto decisionale continuo, con drafting assistito, briefing periodico "
        "e Q&A contestuale sempre disponibili.",
    ),
}

USE_CASE_TITLE = "TCF + One Fiscale: Obiettivo & Funzionamento"

USE_CASE_OBJECTIVE = (
    "Portare l’intelligenza di One Fiscale dentro il Tax Control Framework: sei assistenti "
    "AI dedicati che coprono l’intero ciclo di vita del TCF – dal design della RCM al "
    "monitoraggio continuo delle novità normative – mantenendo i dati cliente nel suo "
    "ambiente Azure e la banca dati fiscale nell’ecosistema WK."
)

USE_CASE_HOW_LINES = [
    "Sei agenti dedicati",
    "Gap Analysis, Driver Assistant, Impact Analysis, Drafting Assistant, Q&A Contestuale e Periodic Briefing – un agente per caso d’uso, orchestrato su FAB.",
    " Doppio indice, dati separati",
    "Indice One Fiscale (banca dati WK) + indice TCF cliente (documenti, strategie, pareri): dati cliente restano nel suo tenant e non si mescolano con altri.",
    " Integrazione via API dedicate",
    "Ogni caso d’uso è esposto come API a partire dal connettore FAB, con data model I/O standardizzato e autenticazione One ID Legal.",
    " Calcoli deterministici affidabili",
    "Le formule dei driver (media pesata su basso/medio/alto) sono eseguite lato TCF: gli agenti forniscono ragionamento e suggerimenti, il TCF resta autorevole sui numeri.",
    " Compliance by design",
    "Data residency EU, GDPR e AI Act coperti dall’ambiente FAB Legal WK; nessuna condivisione della strategia fiscale del cliente al di fuori del perimetro autorizzato.",
]

ARCH_TITLE = "Soluzione Proposta – Integrazione TCF ↔ FAB One Fiscale"

ARCH_BODY = [
    "Panoramica dell’Integrazione",
    "L’applicazione TCF resta ospitata nel suo Azure Cloud dedicato, in un perimetro privato e segregato per il cliente.",
    "Il connettore di integrazione instaura una connessione sicura (HTTPS + mTLS su Private Endpoint) verso la workspace One Fiscale ospitata nel WK Legal Environment su FAB.",
    "Un Orchestrator Agent su FAB instrada ogni richiesta al sub-agent AI competente (uno per caso d’uso) con system prompt dedicato; i calcoli deterministici restano lato TCF.",
    "Due indici separati preservano la governance del dato: banca dati One Fiscale (WK) e documenti privati del cliente TCF (strategie, pareri, big four).",
    "La risposta arricchita viene restituita al TCF che la contestualizza nel wizard, nella RCM o nella fattispecie interpretativa in lavorazione.",
]

IDENTITY_TITLE = "Accesso & Governance dell’Identità"
IDENTITY_SUB = "Un unico modello di identità coerente tra TCF e ambiente WK Legal"

IDENTITY_HEAD_COPY = {
    ("head1", "copy1"): (
        "Modello di Identità",
        [
            "L’accesso ai servizi One Fiscale è basato su utenze One ID Legal come identità di riferimento, coerentemente con lo stack di autenticazione già adottato in ambito WK Legal.",
            "Il TCF propaga l’identità dell’utente all’Orchestrator Agent tramite token firmato, mantenendo la tracciabilità puntuale di ogni chiamata AI.",
        ],
    ),
    ("head2", "copy2"): (
        "Razionale della Scelta",
        [
            "Riutilizzare One ID Legal evita integrazioni SSO custom e centralizza la gestione delle abilitazioni sui servizi One Fiscale.",
            "Il cliente TCF beneficia di onboarding rapido, revoca immediata e audit completo delle interazioni con gli assistenti AI.",
        ],
    ),
    ("head3", "copy3"): (
        "Prerequisiti di Accesso",
        [
            "Ogni utente TCF che utilizza le funzionalità AI deve essere censito in One ID Legal con profilo attivo e abilitato alla workspace One Fiscale.",
            "È richiesta una sottoscrizione a One Fiscale coerente con i casi d’uso previsti (banca dati, fattispecie interpretative, novità normative).",
        ],
    ),
}

POC_TITLE = "Proposta di Approccio – Proof of Concept (PoC)"

POC_HEADER = (
    "Percorso in due fasi per validare in modalità backend l’integrazione tra "
    "il TCF e i sei assistenti AI su FAB One Fiscale"
)

POC_PREREQ = [
    "Prerequisito di Avvio",
    "Il cliente TCF fornisce un set rappresentativo di casi reali: un sotto-insieme di RCM, "
    "una selezione di fattispecie interpretative in corso e un elenco di novità normative recenti.",
    "Questo set consente di validare in ambiente backend il comportamento combinato TCF + FAB, "
    "nel pieno rispetto dei requisiti di sicurezza e riservatezza.",
]

POC_PHASE1 = [
    "Fase 1 – Setup Integrazione & Sei Agenti su FAB",
    "Configurazione della workspace FAB dedicata al progetto TCF e creazione dei sei agenti (Gap Analysis, Driver Assistant, Impact Analysis, Drafting, Q&A Contestuale, Periodic Briefing) con system prompt e sub-agent per indice.",
    "Predisposizione del connettore TCF ↔ FAB (One ID Legal auth, standard I/O JSON, Private Endpoint) e caricamento del set di documenti cliente sull’indice TCF privato.",
    "Output: sei API AI attive e richiamabili dal backend TCF sul set di casi d’uso di prova.",
]

POC_PHASE2 = [
    "Fase 2 – Test end-to-end & Misura Soddisfazione",
    "Esecuzione backend del set di prova su tutti i sei casi d’uso, con generazione di risposte combinate (One Fiscale + documenti cliente TCF) e verifica dei calcoli lato TCF per la parte driver.",
    "Raccolta di feedback strutturato con il cliente (accuratezza, rilevanza, actionability) per ciascun assistente e per il flusso end-to-end.",
    "Output: report di valutazione con grado di soddisfazione per caso d’uso e roadmap di industrializzazione.",
]

POC_INVEST = [
    "Durata & Investimento",
    "Durata complessiva: 8 Settimane",
    "Costo totale PoC: €25.000",
]

COST_TITLE = "Modello Costi della Soluzione TCF + One Fiscale"
COST_SUB = "Stima indicativa di alto livello per l’integrazione dei sei assistenti AI"

COST_VALUES = {
    "d1": ["Costo Iniziale Sviluppo TCF", "€120.000 - €140.000"],
    "a2": ["Costo Iniziale Sviluppo Legal", "€20.000 - €25.000"],
    "b3_manut": ["Costo di Manutenzione", "€15.000/anno"],
    "b3_fab": [
        "Costo Workspace FAB Dedicata",
        "300 €/mese +",
        "",
        "  Quota variabile in funzione di utilizzo utente",
    ],
}


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build():
    if not SRC_PPTX.exists():
        raise SystemExit(f"Template not found: {SRC_PPTX}")
    if not ARCH_PNG.exists() or not CAPS_PNG.exists():
        raise SystemExit("Missing diagram PNGs – run scripts/generate_tfc_diagrams.py first")

    # Copy to a temp working file (bypasses OneDrive reparse points) then overwrite the final path.
    tmp_src = Path(os.environ["TEMP"]) / "uc_template.pptx"
    if not tmp_src.exists() or tmp_src.stat().st_size != SRC_PPTX.stat().st_size:
        shutil.copyfile(SRC_PPTX, tmp_src)
    tmp_out = Path(os.environ["TEMP"]) / "tfc_out.pptx"
    shutil.copyfile(tmp_src, tmp_out)

    prs = Presentation(tmp_out)

    # --- Slide 1 – Cover ---
    s = prs.slides[0]
    set_text_preserving_format(find_shape(s, "Title").text_frame, COVER_TITLE)
    set_text_preserving_format(find_shape(s, "sub").text_frame, COVER_SUB)

    # --- Slide 2 – Contesto e Sfida ---
    s = prs.slides[1]
    set_text_preserving_format(find_shape(s, "Title").text_frame, CONTEXT_TITLE)
    set_text_preserving_format(find_shape(s, "c").text_frame, CONTEXT_HEADER)
    for shape_name, (heading, body) in CONTEXT_BLOCKS.items():
        block = find_shape(s, shape_name)
        set_text_preserving_format(block.text_frame, [heading, body])

    # --- Slide 3 – Obiettivo & Funzionamento ---
    s = prs.slides[2]
    set_text_preserving_format(find_shape(s, "Title").text_frame, USE_CASE_TITLE)
    table = find_shape(s, "Table 7").table
    # Row 0 col 0 = "Obiettivo" (keep the header)
    set_table_cell(table, 1, 0, USE_CASE_OBJECTIVE)
    # Row 2 col 0 = "Come Funziona" (keep the header)
    set_table_cell(table, 3, 0, USE_CASE_HOW_LINES)
    replace_picture(s, "Picture 25", CAPS_PNG)

    # --- Slide 4 – Architettura integrazione ---
    s = prs.slides[3]
    set_text_preserving_format(find_shape(s, "Title").text_frame, ARCH_TITLE)
    set_text_preserving_format(find_shape(s, "Content Placeholder 1").text_frame, ARCH_BODY)
    replace_picture(s, "Picture 3", ARCH_PNG)

    # --- Slide 5 – Identity ---
    s = prs.slides[4]
    set_text_preserving_format(find_shape(s, "Title").text_frame, IDENTITY_TITLE)
    set_text_preserving_format(find_shape(s, "sub").text_frame, IDENTITY_SUB)
    for (head_name, copy_name), (h, body) in IDENTITY_HEAD_COPY.items():
        set_head_copy(s, head_name, copy_name, h, body)

    # --- Slide 6 – PoC ---
    s = prs.slides[5]
    set_text_preserving_format(find_shape(s, "Title").text_frame, POC_TITLE)
    set_text_preserving_format(find_shape(s, "c").text_frame, POC_HEADER)
    set_text_preserving_format(find_shape(s, "TextBox 4").text_frame, POC_PREREQ)
    set_text_preserving_format(find_shape(s, "TextBox 1").text_frame, POC_PHASE1)
    set_text_preserving_format(find_shape(s, "TextBox 3").text_frame, POC_PHASE2)
    set_text_preserving_format(find_shape(s, "TextBox 6").text_frame, POC_INVEST)

    # --- Slide 7 – Cost model ---
    s = prs.slides[6]
    set_text_preserving_format(find_shape(s, "Title").text_frame, COST_TITLE)
    set_text_preserving_format(find_shape(s, "sub").text_frame, COST_SUB)
    # Update named shapes; two shapes are both named 'b3' so map by order.
    b3_shapes = [shp for shp in s.shapes if shp.name == "b3"]
    if len(b3_shapes) >= 1:
        set_text_preserving_format(b3_shapes[0].text_frame, COST_VALUES["b3_manut"])
    if len(b3_shapes) >= 2:
        set_text_preserving_format(b3_shapes[1].text_frame, COST_VALUES["b3_fab"])
    set_text_preserving_format(find_shape(s, "d1").text_frame, COST_VALUES["d1"])
    set_text_preserving_format(find_shape(s, "a2").text_frame, COST_VALUES["a2"])

    # --- Slide 8 – End page: keep as-is ---

    prs.save(tmp_out)
    shutil.copyfile(tmp_out, OUT_PPTX)
    print(f"Wrote: {OUT_PPTX}")


if __name__ == "__main__":
    build()
