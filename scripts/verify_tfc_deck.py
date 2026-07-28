"""Verify the generated TFC deck opens and print each slide's textual content."""
import os, shutil
from pathlib import Path
from pptx import Presentation

src = Path(r"c:\Users\Cesare\OneDrive - NTT DATA EMEAL\Second Brain\reports\Wolters Kluwer - TCF - One Fiscale Integration.pptx")
tmp = Path(os.environ["TEMP"]) / "tfc_verify.pptx"
shutil.copyfile(src, tmp)
prs = Presentation(tmp)

print(f"Slide count: {len(prs.slides)}")
for i, slide in enumerate(prs.slides, 1):
    print(f"\n===== SLIDE {i} (layout: {slide.slide_layout.name}) =====")
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                print(f"[{shape.name}]")
                for line in txt.splitlines():
                    print(f"    {line}")
        elif shape.has_table:
            for r, row in enumerate(shape.table.rows):
                for c, cell in enumerate(row.cells):
                    t = cell.text_frame.text.strip()
                    if t:
                        print(f"[TABLE {r},{c}]")
                        for line in t.splitlines():
                            print(f"    {line}")
        elif shape.shape_type == 13:  # picture
            print(f"[PICTURE {shape.name!r}]  {shape.width}x{shape.height}")
