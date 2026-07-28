"""Inspect the empty Template.pptx to understand its slide masters/layouts."""
import os, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

src = Path(r"c:\Users\Cesare\OneDrive - NTT DATA EMEAL\Second Brain\reports\Template.pptx")
tmp = Path(os.environ["TEMP"]) / "template_probe.pptx"
shutil.copyfile(src, tmp)
prs = Presentation(tmp)

print(f"Slide size: {Emu(prs.slide_width).inches:.2f}in x {Emu(prs.slide_height).inches:.2f}in")
print(f"Slides in template: {len(prs.slides)}")
print(f"Slide masters: {len(prs.slide_masters)}")
for mi, master in enumerate(prs.slide_masters):
    print(f"\nMaster {mi}: {len(master.slide_layouts)} layouts")
    for li, layout in enumerate(master.slide_layouts):
        print(f"  Layout {li}: {layout.name!r}")
        for ph in layout.placeholders:
            print(f"     ph idx={ph.placeholder_format.idx} type={ph.placeholder_format.type} name={ph.name!r} "
                  f"pos=({Emu(ph.left).inches if ph.left else '-'},{Emu(ph.top).inches if ph.top else '-'}) "
                  f"size=({Emu(ph.width).inches if ph.width else '-'},{Emu(ph.height).inches if ph.height else '-'})")

print("\n=== Existing slides ===")
for i, slide in enumerate(prs.slides, 1):
    print(f"\n--- Slide {i} (layout: {slide.slide_layout.name!r}) ---")
    for shape in slide.shapes:
        info = f"  {shape.shape_type} name={shape.name!r}"
        if shape.left is not None:
            info += f" pos=({Emu(shape.left).inches:.2f},{Emu(shape.top).inches:.2f}) size=({Emu(shape.width).inches:.2f},{Emu(shape.height).inches:.2f})"
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                info += f"\n     TEXT: {txt[:200]!r}"
        if shape.is_placeholder:
            info += f"  ph(idx={shape.placeholder_format.idx},type={shape.placeholder_format.type})"
        print(info)
