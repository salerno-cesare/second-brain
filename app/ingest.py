import csv
import html
import json
import re
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

TEXT_EXTENSIONS = {".txt", ".md", ".rst", ".log", ".csv", ".json", ".html", ".htm", ".vtt"}
OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
VTT_TIMESTAMP_RE = re.compile(
    r"^\s*(?:\d{2}:)?\d{2}:\d{2}\.\d{3}\s+-->\s+(?:\d{2}:)?\d{2}:\d{2}\.\d{3}(?:\s+.*)?$"
)


@dataclass(frozen=True)
class SourceFile:
    name: str
    rel_path: str
    ext: str
    size: int
    updated_at: str


def is_supported_file(path: Path) -> bool:
    return path.suffix.lower() in TEXT_EXTENSIONS or path.suffix.lower() in OFFICE_EXTENSIONS or path.suffix.lower() == ".pdf"


def _clean_cell_text(value) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _read_docx_text(path: Path) -> str:
    document = Document(str(path))
    blocks: list[str] = []

    for paragraph in document.paragraphs:
        text = _clean_cell_text(paragraph.text)
        if text:
            blocks.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [text for text in (_clean_cell_text(cell.text) for cell in row.cells) if text]
            if cells:
                blocks.append(" | ".join(cells))

    return "\n".join(blocks)


def _read_pptx_text(path: Path) -> str:
    presentation = Presentation(str(path))
    blocks: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_blocks: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _clean_cell_text(getattr(shape, "text", ""))
                if text:
                    slide_blocks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [text for text in (_clean_cell_text(cell.text) for cell in row.cells) if text]
                    if cells:
                        slide_blocks.append(" | ".join(cells))

        if slide_blocks:
            blocks.append(f"Slide {index}")
            blocks.extend(slide_blocks)

    return "\n".join(blocks)


def _read_xlsx_text(path: Path) -> str:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    blocks: list[str] = []

    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [text for text in (_clean_cell_text(value) for value in row) if text]
                if values:
                    rows.append(" | ".join(values))

            if rows:
                blocks.append(f"Sheet: {sheet.title}")
                blocks.extend(rows)
    finally:
        workbook.close()

    return "\n".join(blocks)


def _read_vtt_text(path: Path) -> str:
    raw = path.read_text(encoding="utf-8-sig", errors="ignore")
    blocks: list[str] = []

    for block in re.split(r"\n\s*\n", raw.replace("\r\n", "\n").replace("\r", "\n")):
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue

        upper_first = lines[0].upper()
        if upper_first == "WEBVTT" or upper_first.startswith(("WEBVTT ", "NOTE", "STYLE", "REGION")):
            continue

        timestamp_index = next((index for index, line in enumerate(lines) if VTT_TIMESTAMP_RE.match(line)), None)
        if timestamp_index is None:
            continue

        cue_lines: list[str] = []
        for line in lines[timestamp_index + 1 :]:
            text = re.sub(r"<[^>]+>", "", line)
            text = html.unescape(text).strip()
            if text:
                cue_lines.append(text)

        if cue_lines:
            blocks.append(" ".join(cue_lines))
    return "\n".join(blocks)


def read_text_from_file(path: Path) -> str:
    ext = path.suffix.lower()

    if ext == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext == ".docx":
        return _read_docx_text(path)

    if ext == ".pptx":
        return _read_pptx_text(path)

    if ext == ".xlsx":
        return _read_xlsx_text(path)

    if ext in {".html", ".htm"}:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        no_tags = re.sub(r"<[^>]+>", " ", raw)
        return html.unescape(no_tags)

    if ext == ".json":
        raw = path.read_text(encoding="utf-8", errors="ignore")
        parsed = json.loads(raw)
        return json.dumps(parsed, indent=2, ensure_ascii=True)

    if ext == ".csv":
        rows: list[str] = []
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        return "\n".join(rows)

    if ext == ".vtt":
        return _read_vtt_text(path)

    return path.read_text(encoding="utf-8", errors="ignore")


def normalize_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def list_source_files(raw_dir: Path, source_dir: Path) -> list[SourceFile]:
    if not raw_dir.exists():
        return []

    source_root = source_dir.resolve()
    files: list[SourceFile] = []
    for path in sorted(raw_dir.rglob("*"), key=lambda item: str(item).lower()):
        if not path.is_file() or not is_supported_file(path):
            continue

        stat = path.stat()
        files.append(
            SourceFile(
                name=path.name,
                rel_path=str(path.resolve().relative_to(source_root)).replace("\\", "/"),
                ext=path.suffix.lower(),
                size=stat.st_size,
                updated_at=datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M"),
            )
        )

    return files
