"""Build TCF / One Fiscale deck from scratch on Template.pptx by drawing every
element with native pptx shapes for full layout control.

Slide list:
  1. Cover                (Master 0 – Innovation Curve on Right)
  2. Contesto e Sfida     (Title Only + 5 native cards)
  3. Obiettivo & Casi d'Uso overview (Title Only + native cards)
  4. Architettura di Integrazione (Title Only + native diagram)
  5. Casi d'Uso in Dettaglio (Title, Subtitle + 6 cards)
  6. Accesso & Governance dell'Identità (Title, Subtitle + 3 blocks)
  7. Proposta PoC (Title Only + 2 phase cards + prerequisite/invest strips)
  8. Modello Costi (Title, Subtitle + 4 cost cards)
  9. End page (Title Only + centered thanks)
"""
from __future__ import annotations

import os
import shutil
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

ROOT = Path(__file__).parent.parent
REPORTS = ROOT / "reports"
SRC = REPORTS / "Template.pptx"
OUT = REPORTS / "Wolters Kluwer - TCF - One Fiscale Integration.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x3B, 0x66)
BLUE = RGBColor(0x1F, 0x5F, 0xA8)
BLUE_LIGHT = RGBColor(0xE8, 0xF1, 0xFB)
BLUE_ULTRA_LIGHT = RGBColor(0xF4, 0xF8, 0xFD)
GREEN = RGBColor(0x0F, 0x8A, 0x3A)
GREEN_LIGHT = RGBColor(0xEA, 0xF7, 0xEE)
GREEN_ULTRA_LIGHT = RGBColor(0xF3, 0xFB, 0xF6)
ORANGE = RGBColor(0xE0, 0x7B, 0x00)
GREY = RGBColor(0x5A, 0x66, 0x72)
GREY_LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _set_text(tf, blocks, *, anchor=MSO_ANCHOR.TOP):
    """Fill a text frame with a list of dicts.

    `blocks` is a list of dicts with keys:
        text     – str (required)
        size     – Pt (default 10)
        bold     – bool (default False)
        color    – RGBColor (default BLACK)
        align    – PP_ALIGN (default LEFT)
        space_before – Pt (default 0)
    """
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    # Clear existing paragraphs
    txBody = tf._txBody
    for p in list(txBody.iter(qn("a:p"))):
        p.getparent().remove(p)

    for i, blk in enumerate(blocks):
        p = etree.SubElement(txBody, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        pPr.set("algn", {
            PP_ALIGN.LEFT: "l",
            PP_ALIGN.CENTER: "ctr",
            PP_ALIGN.RIGHT: "r",
        }.get(blk.get("align", PP_ALIGN.LEFT), "l"))
        if blk.get("space_before"):
            sb = etree.SubElement(pPr, qn("a:spcBef"))
            pts = etree.SubElement(sb, qn("a:spcPts"))
            pts.set("val", str(int(blk["space_before"].pt * 100)))
        r = etree.SubElement(p, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"))
        rPr.set("lang", "it-IT")
        rPr.set("sz", str(int(blk.get("size", Pt(10)).pt * 100)))
        rPr.set("b", "1" if blk.get("bold") else "0")
        rPr.set("dirty", "0")
        color = blk.get("color", BLACK)
        fill = etree.SubElement(rPr, qn("a:solidFill"))
        srgb = etree.SubElement(fill, qn("a:srgbClr"))
        srgb.set("val", "%02X%02X%02X" % (color[0], color[1], color[2]))
        latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", "Calibri")
        t = etree.SubElement(r, qn("a:t"))
        t.text = blk["text"]


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=None, line_w=0.75,
             rounded=False, corner=0.05, shadow=False):
    # EMU coordinates must be integers; callers often use float arithmetic.
    x, y, w, h = int(x), int(y), int(w), int(h)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if rounded:
        # Set corner rounding via adjustment value (0-1 in EMU relative)
        try:
            shp.adjustments[0] = corner
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if not shadow:
        _remove_shadow(shp)
    shp.text_frame.text = ""
    return shp


def _remove_shadow(shape):
    spPr = shape.fill._xPr.find(qn("p:spPr")) if hasattr(shape.fill, "_xPr") else None
    # Simpler: append <a:effectLst/> to spPr to clear inherited shadow
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for eff in spPr.findall(qn("a:effectLst")):
        spPr.remove(eff)
    etree.SubElement(spPr, qn("a:effectLst"))


def add_text_box(slide, x, y, w, h, blocks, *, anchor=MSO_ANCHOR.TOP):
    x, y, w, h = int(x), int(y), int(w), int(h)
    box = slide.shapes.add_textbox(x, y, w, h)
    _set_text(box.text_frame, blocks, anchor=anchor)
    return box


def add_line(slide, x1, y1, x2, y2, *, color=BLUE, weight=1.5,
             arrow_head=True, arrow_tail=False, dashed=False):
    # Normalize direction so python-pptx never emits flipH/flipV, which combined
    # with headEnd/tailEnd triggers PowerPoint's "cannot open" validation.
    if x2 < x1 or (x1 == x2 and y2 < y1):
        x1, x2 = x2, x1
        y1, y2 = y2, y1
        arrow_head, arrow_tail = arrow_tail, arrow_head
    # EMU coordinates must be integers; float arithmetic in the callers must be
    # coerced or PowerPoint refuses to open the file (invalid ST_Coordinate).
    x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    # Ensure children obey OOXML order: fill -> prstDash -> headEnd -> tailEnd.
    if dashed:
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    if arrow_tail:
        head = etree.SubElement(ln, qn("a:headEnd"))
        head.set("type", "triangle")
        head.set("w", "med")
        head.set("len", "med")
    if arrow_head:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return conn


def add_title(slide, title, subtitle=None):
    # Title bar at top
    add_text_box(slide, Inches(0.42), Inches(0.33), Inches(12.5), Inches(0.55),
                 [{"text": title, "size": Pt(24), "bold": True, "color": NAVY}])
    if subtitle:
        add_text_box(slide, Inches(0.42), Inches(0.92), Inches(12.5), Inches(0.42),
                     [{"text": subtitle, "size": Pt(13), "color": GREY}])
    # Thin accent underline
    add_rect(slide, Inches(0.42), Inches(0.90), Inches(1.2), Inches(0.04),
             fill=BLUE)


def add_footer(slide, page_num, total_pages=8):
    # Copyright bottom-left
    add_text_box(slide, Inches(0.42), Inches(7.05), Inches(4.0), Inches(0.3),
                 [{"text": "© 2026 NTT DATA, Inc.", "size": Pt(9), "color": GREY}])
    # Page number center
    add_text_box(slide, Inches(6.0), Inches(7.05), Inches(1.5), Inches(0.3),
                 [{"text": f"{page_num} / {total_pages}", "size": Pt(9),
                   "color": GREY, "align": PP_ALIGN.CENTER}])
    # NOTE: the template master already renders the official NTT DATA logo in
    # the bottom-right corner, so we intentionally do not overlay a text
    # wordmark here. Cover (slide 1) and closing (slide 9) draw their own
    # wordmarks since their layouts do not include the master footer.
    # Bottom accent bar
    add_rect(slide, Inches(0.42), Inches(7.35), Inches(12.5), Inches(0.02),
             fill=BLUE_LIGHT)


def add_ntt_wordmark(slide, x, y, w, h):
    """Render a compact 'NTT DATA' wordmark (as text) in the given box."""
    add_text_box(slide, x, y, w, h,
                 [{"text": "NTT DATA", "size": Pt(14), "bold": True,
                   "color": NAVY, "align": PP_ALIGN.RIGHT}],
                 anchor=MSO_ANCHOR.MIDDLE)


# Master logo used by the Template.pptx (white "N" mark with transparent bg,
# extracted from ppt/media/image1.png). Ships in reports/assets/.
NTT_LOGO_PATH = REPORTS / "assets" / "ntt_master1.png"


def add_ntt_logo_on_dark(slide, right_x, bottom_y, *, mark_h=Inches(0.55)):
    """Place the '[N mark]  NTT DATA' logo (white) with its bottom-right corner
    anchored at (right_x, bottom_y). For use on dark (navy/blue) backgrounds
    such as the cover and end slides."""
    # N mark image (square, transparent WHITE)
    mark_w = mark_h  # image is ~square (447x450)
    text_w = Inches(1.20)
    gap = Inches(0.10)
    total_w = mark_w + gap + text_w
    x0 = right_x - total_w
    y0 = bottom_y - mark_h
    slide.shapes.add_picture(str(NTT_LOGO_PATH), int(x0), int(y0),
                              int(mark_w), int(mark_h))
    add_text_box(slide, x0 + mark_w + gap, y0,
                 text_w, mark_h,
                 [{"text": "NTT DATA", "size": Pt(16), "bold": True,
                   "color": WHITE, "align": PP_ALIGN.LEFT}],
                 anchor=MSO_ANCHOR.MIDDLE)


# Card helpers -------------------------------------------------------------
def card(slide, x, y, w, h, *, title, body_lines, icon_text=None,
         accent=BLUE, bg=WHITE, title_bg=None, title_color=WHITE):
    """A rounded card with an accent header and body text."""
    add_rect(slide, x, y, w, h, fill=bg, line=accent, line_w=1.0,
             rounded=True, corner=0.06)
    header_h = Inches(0.5)
    header_bg = title_bg if title_bg is not None else accent
    add_rect(slide, x, y, w, header_h, fill=header_bg,
             line=header_bg, line_w=0.5, rounded=True, corner=0.12)
    # Cover the bottom rounded corners of the header so it looks like a strip
    add_rect(slide, x, y + Inches(0.15), w, header_h - Inches(0.15),
             fill=header_bg, line=header_bg, line_w=0.5)
    # Title text
    title_blocks = [{"text": title, "size": Pt(11.5), "bold": True,
                     "color": title_color, "align": PP_ALIGN.LEFT}]
    if icon_text:
        title_blocks[0]["text"] = f"{icon_text}  {title}"
    add_text_box(slide, x + Inches(0.05), y + Inches(0.05), w - Inches(0.1),
                 header_h - Inches(0.1), title_blocks, anchor=MSO_ANCHOR.MIDDLE)
    # Body
    body_blocks = []
    for i, line in enumerate(body_lines):
        body_blocks.append({
            "text": line,
            "size": Pt(9.5),
            "color": BLACK,
            "space_before": Pt(4) if i > 0 else Pt(0),
        })
    add_text_box(slide, x + Inches(0.15), y + header_h + Inches(0.1),
                 w - Inches(0.3), h - header_h - Inches(0.2),
                 body_blocks, anchor=MSO_ANCHOR.TOP)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def clear_slides(prs):
    """Remove all existing slides from the presentation."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def slide_title_only(prs):
    # Master 1 (index 1) Layout 0 = "Title Only"
    layout = prs.slide_masters[1].slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    # Remove auto-populated placeholders to have a clean canvas
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)
    # Fill background white
    return slide


def build_slide_1_cover(prs):
    slide = slide_title_only(prs)
    # Full-height navy band on the right (clean, edge-aligned)
    band_x = Inches(9.2)
    band_w = SLIDE_W - band_x
    add_rect(slide, band_x, Inches(0), band_w, SLIDE_H, fill=NAVY)
    # Thin lighter-blue vertical accent strip inside the navy band
    add_rect(slide, band_x, Inches(0), Inches(0.18), SLIDE_H, fill=BLUE)

    # Left accent bar next to the title
    add_rect(slide, Inches(0.42), Inches(2.6), Inches(0.15), Inches(2.6),
             fill=BLUE)

    # Kicker label
    add_text_box(slide, Inches(0.85), Inches(2.1), Inches(7), Inches(0.35),
                 [{"text": "PROPOSTA DI INTEGRAZIONE  •  LUGLIO 2026",
                   "size": Pt(11), "bold": True, "color": BLUE}])

    # Title
    add_text_box(slide, Inches(0.85), Inches(2.5), Inches(8.0), Inches(1.5),
                 [{"text": "Integrazione One Fiscale su TCF",
                   "size": Pt(38), "bold": True, "color": NAVY}])
    # Subtitle
    add_text_box(slide, Inches(0.85), Inches(4.0), Inches(8.0), Inches(2.2),
                 [{"text": "Estensione del Tax Control Framework con capacità AI di "
                           "Wolters Kluwer One Fiscale: sei assistenti dedicati, "
                           "orchestrati su FAB, per gap analysis, driver assistant, "
                           "impact analysis, drafting, Q&A contestuale e briefing "
                           "periodici direttamente nell’ambiente TCF del cliente.",
                   "size": Pt(14), "color": BLACK}])

    # Footer partners (bottom-left)
    add_text_box(slide, Inches(0.85), Inches(6.6), Inches(8), Inches(0.35),
                 [{"text": "Wolters Kluwer  ·  NTT DATA",
                   "size": Pt(12), "bold": True, "color": GREY}])
    add_text_box(slide, Inches(0.85), Inches(6.95), Inches(8), Inches(0.3),
                 [{"text": "© 2026 NTT DATA, Inc.",
                   "size": Pt(9), "color": GREY}])
    # NTT DATA logo on the navy band, bottom-right (white on navy)
    add_ntt_logo_on_dark(slide, right_x=SLIDE_W - Inches(0.30),
                         bottom_y=Inches(7.20))


def build_slide_2_context(prs):
    slide = slide_title_only(prs)
    add_title(slide, "Contesto e Sfida",
              "Il Tax Control Framework deve evolvere per stare al passo con normativa e operatività fiscale")

    blocks = [
        ("Framework fiscale sempre più articolato",
         "Il TCF combina wizard secondo Ag. Entrate, RCM, fattispecie interpretative "
         "e strategia fiscale: mantenere allineamento e completezza è oneroso senza "
         "un supporto intelligente."),
        ("Rischi non mappati e valutazioni soggettive",
         "La configurazione di driver di probabilità e impatto è oggi manuale e "
         "discrezionale, esponendo il cliente a rischi non intercettati e a "
         "valutazioni disomogenee tra periodi."),
        ("Novità normative da tradurre in azioni",
         "Ogni circolare o articolo di One Fiscale può impattare la RCM: serve un "
         "meccanismo automatico che colleghi la novità normativa ai rischi specifici "
         "del cliente."),
        ("Valorizzare la banca dati One Fiscale",
         "One Fiscale è una base normativa qualificata e continuamente aggiornata: "
         "integrarla nel TCF rende ogni funzionalità immediatamente contestualizzata."),
        ("Opportunità di trasformazione operativa",
         "Sei agenti AI dedicati abilitano il salto da compilazione documentale a "
         "supporto decisionale continuo, con drafting, briefing periodico e Q&A "
         "contestuale sempre disponibili."),
    ]
    # 5 cards in 2 rows: 2 opportunity (green, top) + 3 challenge (blue, bottom)
    # Reordered so top row is more prominent and layout is balanced.
    ordered = [blocks[3], blocks[4], blocks[0], blocks[1], blocks[2]]
    accents = [GREEN, GREEN, BLUE, BLUE, BLUE]

    card_w = Inches(4.0)
    card_h = Inches(2.5)
    gap = Inches(0.15)
    top_y = Inches(1.5)
    bot_y = top_y + card_h + gap

    # Top row: 2 green cards, centered
    top_left = Inches(0.42) + (card_w + gap) / 2
    for i in range(2):
        card(slide, top_left + i * (card_w + gap), top_y, card_w, card_h,
             title=ordered[i][0], body_lines=[ordered[i][1]], accent=accents[i])
    # Bottom row: 3 blue cards, full width
    left_start = Inches(0.42)
    for i in range(3):
        card(slide, left_start + i * (card_w + gap), bot_y, card_w, card_h,
             title=ordered[2 + i][0], body_lines=[ordered[2 + i][1]],
             accent=accents[2 + i])

    add_footer(slide, 2)


def build_slide_3_objective(prs):
    slide = slide_title_only(prs)
    add_title(slide, "Obiettivo & Casi d’Uso in Dettaglio",
              "Sei assistenti AI dedicati che coprono l’intero ciclo di vita del TCF — input, cosa fanno, output")

    # --- Top objective banner (full width) ---
    ob_x = Inches(0.42)
    ob_y = Inches(1.45)
    ob_w = Inches(12.5)
    ob_h = Inches(0.9)
    add_rect(slide, ob_x, ob_y, ob_w, ob_h, fill=NAVY, rounded=True, corner=0.04)
    add_text_box(slide, ob_x + Inches(0.30), ob_y + Inches(0.08),
                 Inches(2.0), Inches(0.32),
                 [{"text": "OBIETTIVO", "size": Pt(10.5), "bold": True, "color": BLUE_LIGHT}])
    add_text_box(slide, ob_x + Inches(0.30), ob_y + Inches(0.38),
                 Inches(6.8), Inches(0.50),
                 [{"text": "Portare l’intelligenza di One Fiscale dentro il TCF, mantenendo "
                           "i dati cliente in Azure e la banca dati fiscale nell’ecosistema WK.",
                   "size": Pt(11), "bold": True, "color": WHITE}])
    # Right side chips: key pillars
    chips = ["6 agenti AI", "Doppio indice", "API · One ID Legal", "EU · GDPR · AI Act"]
    chip_area_x = ob_x + Inches(7.35)
    chip_area_w = ob_w - Inches(7.65)
    chip_gap = Inches(0.10)
    chip_w = (chip_area_w - 3 * chip_gap) / 4
    chip_h = Inches(0.42)
    chip_y = ob_y + (ob_h - chip_h) / 2
    for i, label in enumerate(chips):
        cx = chip_area_x + i * (chip_w + chip_gap)
        add_rect(slide, cx, chip_y, chip_w, chip_h, fill=BLUE, rounded=True, corner=0.5)
        add_text_box(slide, cx, chip_y, chip_w, chip_h,
                     [{"text": label, "size": Pt(9), "bold": True, "color": WHITE,
                       "align": PP_ALIGN.CENTER}],
                     anchor=MSO_ANCHOR.MIDDLE)

    # --- 6 use case detail cards (2 x 3 grid) ---
    use_cases = [
        ("1", "Gap Analysis",
         "Confronta la RCM cliente con le linee guida Ag. Entrate (art. 25/11) e segnala rischi non mappati.",
         "Wizard di compilazione TCF",
         "Elenco strutturato di gap con riferimento normativo"),
        ("2", "Driver Assistant",
         "Suggerisce driver, pesi e valori (basso/medio/alto) per probabilità e impatto coerenti con la metodologia cliente.",
         "Config. rischio + Tax Compliance Model",
         "Set di driver proposti, con razionale motivato"),
        ("3", "Impact Analysis",
         "Individua i rischi della RCM impattati da nuove norme, circolari e prassi e propone azioni di aggiornamento.",
         "Novità normative (One Fiscale · circolari)",
         "Rischi impattati + delta di configurazione"),
        ("4", "Drafting Assistant",
         "Propone driver, soglie di materialità e verifiche qualitative/quantitative per compilare la fattispecie.",
         "Fattispecie in lavorazione (es. fusione cross-border)",
         "Bozza strutturata con giustificazione normativa"),
        ("5", "Q&A Contestuale",
         "Risponde su corpus normativo One Fiscale e sul contesto cliente (documenti caricati), con citazioni tracciabili.",
         "Domanda libera dell’utente TCF",
         "Risposta discorsiva con fonti tracciabili"),
        ("6", "Periodic Briefing",
         "Estrae mensilmente le novità fiscali rilevanti filtrate sui processi e rischi specifici del cliente.",
         "Config. cliente (processi · rischi · sottoscrizioni)",
         "Report PDF con highlight e approfondimenti"),
    ]
    grid_x = Inches(0.42)
    grid_y = Inches(2.50)
    cw = Inches(4.05)
    ch = Inches(2.05)
    gx = Inches(0.19)
    gy = Inches(0.12)
    for i, (num, name, do, inp, out) in enumerate(use_cases):
        col = i % 3
        row = i // 3
        x = grid_x + col * (cw + gx)
        y = grid_y + row * (ch + gy)
        # Card
        add_rect(slide, x, y, cw, ch, fill=WHITE, line=BLUE, line_w=0.9,
                 rounded=True, corner=0.05)
        # Header strip (BLUE)
        hdr_h = Inches(0.44)
        add_rect(slide, x, y, cw, hdr_h, fill=BLUE)
        # Number badge (white circle) inside header, left
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       int(x + Inches(0.10)),
                                       int(y + Inches(0.07)),
                                       int(Inches(0.30)),
                                       int(Inches(0.30)))
        badge.fill.solid(); badge.fill.fore_color.rgb = WHITE
        badge.line.fill.background(); _remove_shadow(badge)
        _set_text(badge.text_frame,
                  [{"text": num, "size": Pt(11), "bold": True, "color": BLUE,
                    "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        # Name
        add_text_box(slide, x + Inches(0.50), y, cw - Inches(0.60), hdr_h,
                     [{"text": name, "size": Pt(12), "bold": True, "color": WHITE}],
                     anchor=MSO_ANCHOR.MIDDLE)
        # Body: "cosa fa"
        body_y = y + hdr_h + Inches(0.06)
        body_h = Inches(0.78)
        add_text_box(slide, x + Inches(0.15), body_y, cw - Inches(0.30), body_h,
                     [{"text": do, "size": Pt(9.5), "color": BLACK}])
        # Input row
        io_y = body_y + body_h
        io_h = Inches(0.36)
        add_text_box(slide, x + Inches(0.15), io_y, Inches(0.55), io_h,
                     [{"text": "IN", "size": Pt(8.5), "bold": True, "color": BLUE,
                       "align": PP_ALIGN.CENTER}],
                     anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, x + Inches(0.72), io_y, cw - Inches(0.87), io_h,
                     [{"text": inp, "size": Pt(8.5), "color": GREY}],
                     anchor=MSO_ANCHOR.MIDDLE)
        # Output row
        out_y = io_y + io_h
        out_h = Inches(0.36)
        add_text_box(slide, x + Inches(0.15), out_y, Inches(0.55), out_h,
                     [{"text": "OUT", "size": Pt(8.5), "bold": True, "color": GREEN,
                       "align": PP_ALIGN.CENTER}],
                     anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, x + Inches(0.72), out_y, cw - Inches(0.87), out_h,
                     [{"text": out, "size": Pt(8.5), "color": GREY}],
                     anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 3)


def build_slide_4_architecture(prs):
    slide = slide_title_only(prs)
    add_title(slide, "Architettura di Integrazione",
              "TCF Application (Azure) ↔ Workspace One Fiscale (FAB · WK Legal) — ingestion & trigger schedulati")

    # ------------------------------------------------------------------
    # LAYOUT PLAN (all coords in Inches, both clouds share row Y-levels)
    #   Row A (data)     y = 2.20 – 3.10  → DB / Storage        vs   Indici
    #   Row B (compute)  y = 3.30 – 4.20  → TCF App / Connector vs   Orchestrator
    #   Row C (bottom)   y = 4.60 – 5.70  → Scheduler (wide)    vs   6 agenti (single row)
    # ------------------------------------------------------------------

    # Cloud containers (height sized so the bottom padding sits above the legend)
    az_x, az_y = Inches(0.42), Inches(1.55)
    az_w, az_h = Inches(5.5), Inches(4.60)
    wk_x, wk_y = Inches(6.95), Inches(1.55)
    wk_w, wk_h = Inches(5.95), Inches(4.60)

    add_rect(slide, az_x, az_y, az_w, az_h,
             fill=BLUE_ULTRA_LIGHT, line=BLUE, line_w=1.5, rounded=True, corner=0.03)
    add_rect(slide, wk_x, wk_y, wk_w, wk_h,
             fill=GREEN_ULTRA_LIGHT, line=GREEN, line_w=1.5, rounded=True, corner=0.03)

    # ---- Nested helpers -----------------------------------------------
    def _cloud_header(x, y, w, h, text, color):
        add_rect(slide, x, y, w, h, fill=color)
        add_text_box(slide, x, y, w, h,
                     [{"text": text, "size": Pt(11), "bold": True,
                       "color": WHITE, "align": PP_ALIGN.CENTER}],
                     anchor=MSO_ANCHOR.MIDDLE)

    def _tile(x, y, w, h, *, title, desc, border, fill=WHITE,
              title_color=None, desc_color=BLACK):
        title_color = title_color or border
        add_rect(slide, x, y, w, h, fill=fill, line=border, line_w=1.1,
                 rounded=True, corner=0.06)
        add_text_box(slide, x, y + Inches(0.10), w, Inches(0.34),
                     [{"text": title, "size": Pt(11), "bold": True,
                       "color": title_color, "align": PP_ALIGN.CENTER}])
        add_text_box(slide, x + Inches(0.12), y + Inches(0.46),
                     w - Inches(0.24), h - Inches(0.52),
                     [{"text": desc, "size": Pt(9), "color": desc_color,
                       "align": PP_ALIGN.CENTER}])

    # ---- Cloud headers -------------------------------------------------
    hdr_h = Inches(0.4)
    _cloud_header(az_x + Inches(0.18), Inches(1.65),
                  az_w - Inches(0.36), hdr_h,
                  "AZURE CLOUD  ·  Cliente TCF", BLUE)
    _cloud_header(wk_x + Inches(0.18), Inches(1.65),
                  wk_w - Inches(0.36), hdr_h,
                  "WK LEGAL ENVIRONMENT  ·  FAB  (Workspace One Fiscale)", GREEN)

    # ---- Row anchors (shared) ------------------------------------------
    row_a_y = Inches(2.20)
    row_a_h = Inches(0.90)
    row_b_y = Inches(3.30)
    row_b_h = Inches(0.90)
    row_c_y = Inches(4.60)   # bottom row: Scheduler + 6 agents
    row_c_h = Inches(1.10)

    # ---- Left cloud (Azure) content ------------------------------------
    inner_x_l = az_x + Inches(0.18)
    inner_w_l = az_w - Inches(0.36)
    box_gap = Inches(0.15)
    box_w_l = (inner_w_l - box_gap) / 2

    # Row A: DB TCF (LEFT), Storage Documentale (RIGHT – near gap for ingestion arrow)
    db_x = inner_x_l
    storage_x = inner_x_l + box_w_l + box_gap
    _tile(db_x, row_a_y, box_w_l, row_a_h,
          title="DB TCF",
          desc="Database applicativo — RCM · driver · wizard",
          border=BLUE)
    _tile(storage_x, row_a_y, box_w_l, row_a_h,
          title="Storage Documentale",
          desc="Documenti privati cifrati del cliente TCF",
          border=BLUE)

    # Row B: TCF Application (LEFT), Integration Connector (RIGHT – near gap for HTTPS arrow)
    tcf_x = inner_x_l
    conn_x = inner_x_l + box_w_l + box_gap
    _tile(tcf_x, row_b_y, box_w_l, row_b_h,
          title="TCF Application",
          desc="Web app cliente — wizard · RCM · fattispecie",
          border=NAVY)
    _tile(conn_x, row_b_y, box_w_l, row_b_h,
          title="Integration Connector",
          desc="Gateway API — JSON · mTLS",
          border=NAVY)

    # Row C: Scheduler Job (full inner width, aligned with the agents row)
    add_rect(slide, inner_x_l, row_c_y, inner_w_l, row_c_h,
             fill=ORANGE, rounded=True, corner=0.05)
    add_text_box(slide, inner_x_l, row_c_y + Inches(0.18),
                 inner_w_l, Inches(0.34),
                 [{"text": "SCHEDULER JOB  (Azure)",
                   "size": Pt(11), "bold": True, "color": WHITE,
                   "align": PP_ALIGN.CENTER}])
    add_text_box(slide, inner_x_l + Inches(0.2), row_c_y + Inches(0.55),
                 inner_w_l - Inches(0.4), row_c_h - Inches(0.60),
                 [{"text": "Job schedulato mensile che invoca il Periodic Briefing Agent per generare il report",
                   "size": Pt(9), "color": WHITE, "align": PP_ALIGN.CENTER}])

    # ---- Right cloud (FAB) content -------------------------------------
    inner_x_r = wk_x + Inches(0.18)
    inner_w_r = wk_w - Inches(0.36)
    idx_w = (inner_w_r - box_gap) / 2

    # Row A: Indice TCF Cliente (LEFT – near gap for ingestion), Indice One Fiscale (RIGHT)
    _tile(inner_x_r, row_a_y, idx_w, row_a_h,
          title="Indice TCF Cliente",
          desc="Documenti vettorizzati (aggiornati via ingestion)",
          border=GREEN)
    _tile(inner_x_r + idx_w + box_gap, row_a_y, idx_w, row_a_h,
          title="Indice One Fiscale",
          desc="Banca dati normativa WK (aggiornata continuamente)",
          border=GREEN)

    # Row B: Orchestrator Agent (wide, NAVY prominent)
    add_rect(slide, inner_x_r, row_b_y, inner_w_r, row_b_h,
             fill=NAVY, rounded=True, corner=0.06)
    add_text_box(slide, inner_x_r, row_b_y + Inches(0.1),
                 inner_w_r, Inches(0.35),
                 [{"text": "ORCHESTRATOR AGENT",
                   "size": Pt(12), "bold": True, "color": WHITE,
                   "align": PP_ALIGN.CENTER}])
    add_text_box(slide, inner_x_r, row_b_y + Inches(0.48),
                 inner_w_r, Inches(0.42),
                 [{"text": "Routing verso i 6 agenti · retrieval dai due indici · policy & audit",
                   "size": Pt(9.5), "color": BLUE_LIGHT,
                   "align": PP_ALIGN.CENTER}])

    # Row C: 6 agents in a single row.
    # Periodic Briefing at col 0 (leftmost) so the Scheduler -> Periodic
    # Briefing horizontal arrow is a straight, non-crossing line.
    ag_gap = Inches(0.08)
    ag_w = (inner_w_r - 5 * ag_gap) / 6
    ag_h = row_c_h

    agents = [
        # (name, is_periodic_briefing)
        ("Periodic\nBriefing",  True),
        ("Gap\nAnalysis",       False),
        ("Driver\nAssistant",   False),
        ("Impact\nAnalysis",    False),
        ("Drafting\nAssistant", False),
        ("Q&A\nContestuale",    False),
    ]
    agent_centers_x = []
    for i, (name, is_pb) in enumerate(agents):
        x = inner_x_r + i * (ag_w + ag_gap)
        agent_centers_x.append(x + ag_w / 2)
        if is_pb:
            add_rect(slide, x, row_c_y, ag_w, ag_h,
                     fill=ORANGE, rounded=True, corner=0.06)
            txt_color = WHITE
        else:
            add_rect(slide, x, row_c_y, ag_w, ag_h,
                     fill=WHITE, line=GREEN, line_w=1.1,
                     rounded=True, corner=0.06)
            txt_color = NAVY
        add_text_box(slide, x, row_c_y, ag_w, ag_h,
                     [{"text": name.replace("\n", " "),
                       "size": Pt(9.5), "bold": True,
                       "color": txt_color, "align": PP_ALIGN.CENTER}],
                     anchor=MSO_ANCHOR.MIDDLE)

    # ---- Internal arrows: Azure ----------------------------------------
    # TCF App -> DB TCF  (vertical UP)
    tcf_top_cx = tcf_x + box_w_l / 2
    add_line(slide, tcf_top_cx, row_b_y,
             tcf_top_cx, row_a_y + row_a_h,
             color=GREY, weight=1.2, arrow_head=True)

    # TCF App -> Integration Connector  (horizontal RIGHT)
    tcf_right_x = tcf_x + box_w_l
    conn_left_x = conn_x
    row_b_cy = row_b_y + row_b_h / 2
    add_line(slide, tcf_right_x, row_b_cy,
             conn_left_x, row_b_cy,
             color=GREY, weight=1.2, arrow_head=True)

    # TCF App -> Storage Documentale  (short diagonal through the corner gap)
    add_line(slide, tcf_right_x, row_b_y,
             storage_x + box_w_l * 0.35, row_a_y + row_a_h,
             color=GREY, weight=1.2, arrow_head=True)

    # ---- Internal arrows: FAB ------------------------------------------
    # Orchestrator <-> 2 indexes (short bidirectional vertical arrows).
    idx_bot_y = row_a_y + row_a_h
    orch_top_y = row_b_y
    for xi in [inner_x_r + idx_w / 2,
               inner_x_r + idx_w + box_gap + idx_w / 2]:
        add_line(slide, xi, idx_bot_y, xi, orch_top_y,
                 color=GREEN, weight=1.1, arrow_head=True, arrow_tail=True)

    # Orchestrator -> each of the 6 agents (single vertical arrow per agent).
    orch_bot_y = row_b_y + row_b_h
    for cx in agent_centers_x:
        add_line(slide, cx, orch_bot_y, cx, row_c_y,
                 color=GREEN, weight=0.9, arrow_head=True)

    # ---- Cross-cloud arrows --------------------------------------------
    # 1. Integration Connector <-> Orchestrator  (HTTPS · mTLS, BLUE bidirectional)
    conn_right_x = conn_x + box_w_l
    conn_center_y = row_b_y + row_b_h / 2
    orch_left_x = inner_x_r
    add_line(slide, conn_right_x, conn_center_y, orch_left_x, conn_center_y,
             color=BLUE, weight=3.0, arrow_head=True, arrow_tail=True)
    chip_w = Inches(0.95)
    chip_h = Inches(0.5)
    chip_cx = (conn_right_x + orch_left_x) / 2
    add_rect(slide, chip_cx - chip_w / 2, conn_center_y - chip_h / 2,
             chip_w, chip_h, fill=WHITE, line=BLUE, line_w=1.0,
             rounded=True, corner=0.3)
    add_text_box(slide, chip_cx - chip_w / 2, conn_center_y - chip_h / 2,
                 chip_w, chip_h,
                 [{"text": "HTTPS\nmTLS", "size": Pt(8.5), "bold": True,
                   "color": BLUE, "align": PP_ALIGN.CENTER}],
                 anchor=MSO_ANCHOR.MIDDLE)

    # 2. Storage Documentale -> Indice TCF Cliente  (Ingestion Job, GREY dashed)
    storage_right_x = storage_x + box_w_l
    storage_row_a_cy = row_a_y + row_a_h / 2
    idx_tcf_left_x = inner_x_r
    add_line(slide, storage_right_x, storage_row_a_cy,
             idx_tcf_left_x, storage_row_a_cy,
             color=GREY, weight=1.6, arrow_head=True, arrow_tail=False,
             dashed=True)
    lbl2_cx = (storage_right_x + idx_tcf_left_x) / 2
    add_text_box(slide, lbl2_cx - Inches(0.65), storage_row_a_cy - Inches(0.38),
                 Inches(1.3), Inches(0.28),
                 [{"text": "Ingestion Job", "size": Pt(8.5), "bold": True,
                   "color": GREY, "align": PP_ALIGN.CENTER}])

    # 3. Scheduler Job -> Periodic Briefing Agent  (Monthly Trigger, ORANGE dashed)
    sched_right_x = inner_x_l + inner_w_l
    sched_center_y = row_c_y + row_c_h / 2
    pb_left_x = inner_x_r  # Periodic Briefing is at col 0 (leftmost)
    add_line(slide, sched_right_x, sched_center_y,
             pb_left_x, sched_center_y,
             color=ORANGE, weight=1.8, arrow_head=True, arrow_tail=False,
             dashed=True)
    lbl3_cx = (sched_right_x + pb_left_x) / 2
    add_text_box(slide, lbl3_cx - Inches(0.75), sched_center_y - Inches(0.38),
                 Inches(1.5), Inches(0.28),
                 [{"text": "Monthly Trigger", "size": Pt(8.5), "bold": True,
                   "color": ORANGE, "align": PP_ALIGN.CENTER}])

    # ---- Bottom legend strip (placed OUTSIDE the clouds) ---------------
    legend_y = Inches(6.35)
    add_text_box(slide, Inches(0.42), legend_y, Inches(12.5), Inches(0.28),
                 [{"text": ("Legenda:   ▬▬  HTTPS · mTLS bidirezionale        "
                            "▬ ▬  Ingestion (batch)        "
                            "▬ ▬  Monthly Trigger (scheduled)"),
                   "size": Pt(8.5), "color": GREY,
                   "align": PP_ALIGN.CENTER}])

    add_footer(slide, 4)


def build_slide_5_use_cases_detail(prs):
    slide = slide_title_only(prs)
    add_title(slide, "Casi d’Uso in Dettaglio",
              "Cosa fa ciascun assistente AI e come si integra nel workflow TCF")

    cases = [
        ("Gap Analysis",
         ["Input", "Wizard di compilazione TCF secondo Ag. Entrate."],
         ["Cosa fa", "Confronta la RCM del cliente con le linee guida (art. 25/11) "
                     "e segnala rischi non mappati o incompleti."],
         ["Output", "Elenco strutturato di gap con riferimento normativo."]),
        ("Driver Assistant",
         ["Input", "Configurazione rischio + Tax Compliance Model cliente."],
         ["Cosa fa", "Suggerisce driver, pesi e valori (basso/medio/alto) per "
                     "probabilità e impatto in coerenza con la metodologia cliente."],
         ["Output", "Set di driver proposti, con razionale motivato."]),
        ("Impact Analysis",
         ["Input", "Novità normative (One Fiscale, circolari, prassi)."],
         ["Cosa fa", "Individua i rischi della RCM impattati dalla nuova norma e "
                     "propone azioni di aggiornamento."],
         ["Output", "Lista rischi impattati + delta di configurazione suggerito."]),
        ("Drafting Assistant",
         ["Input", "Fattispecie interpretativa in lavorazione (es. fusione cross-border)."],
         ["Cosa fa", "Propone driver, soglie di materialità e verifiche "
                     "qualitative/quantitative per compilare la fattispecie."],
         ["Output", "Bozza strutturata di fattispecie con giustificazione normativa."]),
        ("Q&A Contestuale",
         ["Input", "Domanda libera dell’utente TCF."],
         ["Cosa fa", "Risponde su corpus normativo One Fiscale e sul contesto "
                     "cliente (documenti caricati), con citazioni."],
         ["Output", "Risposta discorsiva con fonti tracciabili."]),
        ("Periodic Briefing",
         ["Input", "Configurazione cliente (processi, rischi, sottoscrizioni)."],
         ["Cosa fa", "Estrae mensilmente le novità fiscali rilevanti filtrate su "
                     "processi e rischi specifici."],
         ["Output", "Report PDF scaricabile con highlight e link agli approfondimenti."]),
    ]
    grid_x = Inches(0.42)
    grid_y = Inches(1.55)
    cw = Inches(4.0)
    ch = Inches(2.55)
    gx = Inches(0.15)
    gy = Inches(0.15)
    for i, (name, inp, do, out) in enumerate(cases):
        col = i % 3
        row = i // 3
        x = grid_x + col * (cw + gx)
        y = grid_y + row * (ch + gy)
        add_rect(slide, x, y, cw, ch, fill=WHITE, line=BLUE, line_w=0.9,
                 rounded=True, corner=0.05)
        # Header strip
        add_rect(slide, x, y, cw, Inches(0.45), fill=BLUE)
        add_text_box(slide, x + Inches(0.15), y, cw - Inches(0.3), Inches(0.45),
                     [{"text": name, "size": Pt(12), "bold": True, "color": WHITE}],
                     anchor=MSO_ANCHOR.MIDDLE)
        # Body: 3 mini rows
        row_h = (ch - Inches(0.55)) / 3
        for j, (lab, txt) in enumerate([inp, do, out]):
            ry = y + Inches(0.5) + j * row_h
            add_text_box(slide, x + Inches(0.15), ry, Inches(0.9), row_h,
                         [{"text": lab, "size": Pt(9), "bold": True, "color": BLUE}])
            add_text_box(slide, x + Inches(1.05), ry, cw - Inches(1.2), row_h,
                         [{"text": txt, "size": Pt(9), "color": BLACK}])

    add_footer(slide, 5)


def build_slide_6_identity(prs):
    slide = slide_title_only(prs)
    add_title(slide, "Accesso & Governance dell’Identità",
              "Un unico modello di identità coerente tra TCF e ambiente WK Legal")

    blocks = [
        ("Modello di Identità",
         "L’accesso ai servizi One Fiscale è basato su utenze One ID Legal come "
         "identità di riferimento, coerentemente con lo stack di autenticazione già "
         "adottato in ambito WK Legal.\n\nIl TCF propaga l’identità dell’utente "
         "all’Orchestrator Agent tramite token firmato, mantenendo la tracciabilità "
         "puntuale di ogni chiamata AI."),
        ("Razionale della Scelta",
         "Riutilizzare One ID Legal evita integrazioni SSO custom e centralizza la "
         "gestione delle abilitazioni sui servizi One Fiscale.\n\nIl cliente TCF "
         "beneficia di onboarding rapido, revoca immediata e audit completo delle "
         "interazioni con gli assistenti AI."),
        ("Prerequisiti di Accesso",
         "Ogni utente TCF che utilizza le funzionalità AI deve essere censito in "
         "One ID Legal con profilo attivo e abilitato alla workspace One Fiscale."
         "\n\nÈ richiesta una sottoscrizione a One Fiscale coerente con i casi "
         "d’uso previsti (banca dati, fattispecie, novità normative)."),
    ]

    col_w = Inches(4.0)
    col_gap = Inches(0.2)
    col_y = Inches(1.7)
    col_h = Inches(5.1)
    total_w = col_w * 3 + col_gap * 2
    start_x = (SLIDE_W - total_w) / 2

    for i, (h, body) in enumerate(blocks):
        x = start_x + i * (col_w + col_gap)
        # Card (full height, badge sits inside)
        add_rect(slide, x, col_y, col_w, col_h,
                 fill=BLUE_ULTRA_LIGHT, line=BLUE, line_w=0.8, rounded=True, corner=0.05)
        # Number badge fully inside the card, top-left
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        x + Inches(0.3), col_y + Inches(0.3),
                                        Inches(0.6), Inches(0.6))
        badge.fill.solid(); badge.fill.fore_color.rgb = BLUE
        badge.line.fill.background(); _remove_shadow(badge)
        _set_text(badge.text_frame,
                  [{"text": str(i + 1), "size": Pt(18), "bold": True, "color": WHITE,
                    "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        # Heading next to badge
        add_text_box(slide, x + Inches(1.05), col_y + Inches(0.32),
                     col_w - Inches(1.3), Inches(0.55),
                     [{"text": h, "size": Pt(13), "bold": True, "color": NAVY}],
                     anchor=MSO_ANCHOR.MIDDLE)
        # Divider
        add_rect(slide, x + Inches(0.3), col_y + Inches(1.1),
                 col_w - Inches(0.6), Inches(0.03), fill=BLUE_LIGHT)
        # Body
        # Split body on \n\n for paragraphs
        paras = [p.strip() for p in body.split("\n\n") if p.strip()]
        body_blocks = []
        for j, p in enumerate(paras):
            body_blocks.append({
                "text": p,
                "size": Pt(10),
                "color": BLACK,
                "space_before": Pt(6) if j > 0 else Pt(0),
            })
        add_text_box(slide, x + Inches(0.3), col_y + Inches(1.3),
                     col_w - Inches(0.6), col_h - Inches(1.5),
                     body_blocks)

    add_footer(slide, 5)


def build_slide_7_poc(prs):
    slide = slide_title_only(prs)
    add_title(slide, "Proposta di Approccio – Proof of Concept (PoC)",
              "Percorso in due fasi focalizzato sul solo caso d’uso Gap Analysis per validare l’integrazione TCF ↔ FAB One Fiscale")

    # Prerequisite strip
    pr_x, pr_y = Inches(0.42), Inches(1.7)
    pr_w, pr_h = Inches(12.5), Inches(0.85)
    add_rect(slide, pr_x, pr_y, pr_w, pr_h, fill=GREY_LIGHT, rounded=True, corner=0.06)
    add_text_box(slide, pr_x + Inches(0.3), pr_y + Inches(0.1),
                 Inches(3.5), Inches(0.35),
                 [{"text": "PREREQUISITO", "size": Pt(11), "bold": True, "color": ORANGE}])
    add_text_box(slide, pr_x + Inches(0.3), pr_y + Inches(0.42),
                 pr_w - Inches(0.6), pr_h - Inches(0.5),
                 [{"text": "Il cliente TCF fornisce un set di documenti di esempio "
                           "(RCM, policy interne, linee guida) e un set di domande & risposte "
                           "attese da usare come benchmark di valutazione del caso d’uso Gap Analysis.",
                   "size": Pt(10), "color": BLACK}])

    # Two phase cards
    phase_y = Inches(2.75)
    phase_h = Inches(3.5)
    phase_w = Inches(6.05)
    gap = Inches(0.4)

    def _phase(x, num, title, bullets, output, color):
        add_rect(slide, x, phase_y, phase_w, phase_h, fill=WHITE,
                 line=color, line_w=1.2, rounded=True, corner=0.04)
        # Header strip (colored) that hosts the badge + title
        header_h = Inches(0.85)
        add_rect(slide, x, phase_y, phase_w, header_h,
                 fill=color, rounded=True, corner=0.08)
        add_rect(slide, x, phase_y + Inches(0.35), phase_w, header_h - Inches(0.35),
                 fill=color)
        # Number badge sits inside header (white circle on colored strip)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        x + Inches(0.25),
                                        phase_y + (header_h - Inches(0.55)) / 2,
                                        Inches(0.55), Inches(0.55))
        badge.fill.solid(); badge.fill.fore_color.rgb = WHITE
        badge.line.fill.background(); _remove_shadow(badge)
        _set_text(badge.text_frame,
                  [{"text": str(num), "size": Pt(16), "bold": True, "color": color,
                    "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, x + Inches(1.0), phase_y,
                     phase_w - Inches(1.2), header_h,
                     [{"text": title, "size": Pt(13), "bold": True, "color": WHITE}],
                     anchor=MSO_ANCHOR.MIDDLE)
        # Bullets
        bl_blocks = []
        for i, b in enumerate(bullets):
            bl_blocks.append({
                "text": "•  " + b,
                "size": Pt(10),
                "color": BLACK,
                "space_before": Pt(4) if i > 0 else Pt(0),
            })
        add_text_box(slide, x + Inches(0.3), phase_y + header_h + Inches(0.15),
                     phase_w - Inches(0.6),
                     phase_h - header_h - Inches(0.95),
                     bl_blocks)
        # Output strip
        out_y = phase_y + phase_h - Inches(0.75)
        add_rect(slide, x + Inches(0.3), out_y, phase_w - Inches(0.6), Inches(0.55),
                 fill=BLUE_LIGHT, rounded=True, corner=0.15)
        add_text_box(slide, x + Inches(0.5), out_y, phase_w - Inches(1.0), Inches(0.55),
                     [{"text": f"Output → {output}", "size": Pt(9.5), "bold": True,
                       "color": NAVY}],
                     anchor=MSO_ANCHOR.MIDDLE)

    _phase(Inches(0.42), 1,
           "Fase 1 – Setup Integrazione & Agente Gap Analysis su FAB",
           [
               "Configurazione workspace FAB dedicata al progetto TCF",
               "Creazione del solo agente Gap Analysis con system prompt dedicato",
               "Predisposizione connettore TCF ↔ FAB (One ID Legal, JSON I/O, HTTPS · mTLS)",
               "Caricamento del set di documenti di esempio sull’indice TCF privato",
           ],
           "API Gap Analysis attiva e richiamabile dal backend TCF sul set di test",
           BLUE)

    _phase(Inches(0.42) + phase_w + gap, 2,
           "Fase 2 – Test end-to-end & Confronto Risposte",
           [
               "Esecuzione backend del set di domande sul caso d’uso Gap Analysis",
               "Generazione delle risposte combinando One Fiscale + documenti cliente",
               "Confronto tra le risposte dell’agente e le risposte attese fornite dal cliente",
               "Raccolta feedback strutturato: accuratezza, rilevanza, actionability",
           ],
           "Report di valutazione su Gap Analysis + roadmap di estensione agli altri 5 agenti",
           GREEN)

    # Bottom investment strip
    inv_y = Inches(6.4)
    add_rect(slide, Inches(0.42), inv_y, Inches(12.5), Inches(0.5),
             fill=NAVY, rounded=True, corner=0.15)
    add_text_box(slide, Inches(0.6), inv_y, Inches(12.15), Inches(0.5),
                 [{"text": "Durata complessiva:  8 settimane        "
                           "Costo totale PoC:  €25.000",
                   "size": Pt(12), "bold": True, "color": WHITE,
                   "align": PP_ALIGN.CENTER}],
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, 6)


def build_slide_8_costs(prs):
    slide = slide_title_only(prs)
    add_title(slide, "Modello Costi della Soluzione TCF + One Fiscale",
              "Voci di costo identificate — valori economici da definire in fase di scoping")

    def _cost_card(x, y, w, h, label, price, desc, color):
        add_rect(slide, x, y, w, h, fill=WHITE, line=color, line_w=1.4,
                 rounded=True, corner=0.03)
        # Left color band
        add_rect(slide, x, y, Inches(0.25), h, fill=color, rounded=True, corner=0.2)
        add_rect(slide, x + Inches(0.15), y, Inches(0.1), h, fill=color)
        # Label
        add_text_box(slide, x + Inches(0.45), y + Inches(0.15), w - Inches(0.6), Inches(0.42),
                     [{"text": label, "size": Pt(11.5), "bold": True, "color": GREY}])
        # Price (question-mark placeholder)
        add_text_box(slide, x + Inches(0.45), y + Inches(0.55), w - Inches(0.6), Inches(0.6),
                     [{"text": price, "size": Pt(22), "bold": True, "color": color}])
        # Description
        add_text_box(slide, x + Inches(0.45), y + Inches(1.20), w - Inches(0.6),
                     h - Inches(1.30),
                     [{"text": desc, "size": Pt(9.5), "color": BLACK}])

    # =========================================================
    # Section 1: Costi di Progetto NTT DATA (3 cards)
    # =========================================================
    sec1_y = Inches(1.55)
    add_text_box(slide, Inches(0.42), sec1_y, Inches(12.5), Inches(0.35),
                 [{"text": "COSTI DI PROGETTO NTT DATA",
                   "size": Pt(11), "bold": True, "color": NAVY}])
    add_rect(slide, Inches(0.42), sec1_y + Inches(0.32), Inches(12.5), Inches(0.02),
             fill=NAVY)

    row1_y = Inches(1.98)
    row1_h = Inches(2.10)
    proj_gx = Inches(0.20)
    proj_cw = (Inches(12.5) - 2 * proj_gx) / 3

    project_costs = [
        ("Sviluppo TCF",
         "€ ?",
         "Sviluppo del connettore, adattamento wizard/RCM, integrazione UI e "
         "orchestrazione lato TCF per i casi d’uso.",
         NAVY),
        ("Sviluppo Legal (FAB)",
         "€ ?",
         "Configurazione workspace FAB, creazione degli agenti, ingegneria "
         "dei system prompt e allineamento indici.",
         BLUE),
        ("Manutenzione",
         "€ ? / anno",
         "Manutenzione evolutiva, tuning dei prompt, aggiornamento dataset "
         "e monitoraggio qualità delle risposte.",
         GREEN),
    ]
    for i, (label, price, desc, color) in enumerate(project_costs):
        x = Inches(0.42) + i * (proj_cw + proj_gx)
        _cost_card(x, row1_y, proj_cw, row1_h, label, price, desc, color)

    # =========================================================
    # Section 2: Costi Piattaforma FAB · WK Legal (2 cards)
    # =========================================================
    sec2_y = Inches(4.25)
    add_text_box(slide, Inches(0.42), sec2_y, Inches(12.5), Inches(0.35),
                 [{"text": "COSTI PIATTAFORMA FAB · WK LEGAL",
                   "size": Pt(11), "bold": True, "color": ORANGE}])
    add_rect(slide, Inches(0.42), sec2_y + Inches(0.32), Inches(12.5), Inches(0.02),
             fill=ORANGE)

    row2_y = Inches(4.68)
    row2_h = Inches(1.85)
    fab_gx = Inches(0.4)
    fab_cw = (Inches(12.5) - fab_gx) / 2

    fab_costs = [
        ("Costo Variabile per utilizzo LLM",
         "€ ?",
         "Quota variabile in funzione del volume di chiamate agli agenti / token consumati.",
         ORANGE),
        ("Costo fisso TCF Index FAB",
         "€ ? / mese",
         "Costo fisso di piattaforma per l’indice TCF Cliente dedicato sulla workspace FAB.",
         BLUE),
    ]
    for i, (label, price, desc, color) in enumerate(fab_costs):
        x = Inches(0.42) + i * (fab_cw + fab_gx)
        _cost_card(x, row2_y, fab_cw, row2_h, label, price, desc, color)

    # Note above the footer, clear of the page number
    add_text_box(slide, Inches(0.42), Inches(6.65), Inches(12.5), Inches(0.32),
                 [{"text": "Tutte le cifre sono da definire in fase di scoping "
                           "contrattuale; escluso costi infrastrutturali Azure lato cliente.",
                   "size": Pt(9), "color": GREY, "align": PP_ALIGN.CENTER}])

    add_footer(slide, 7)


def build_slide_9_end(prs):
    slide = slide_title_only(prs)
    # Background bands
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(3.5), fill=NAVY)
    add_rect(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(4.0), fill=BLUE)
    # Title
    add_text_box(slide, Inches(0.42), Inches(2.4), Inches(12.5), Inches(1.4),
                 [{"text": "Grazie", "size": Pt(60), "bold": True, "color": WHITE,
                   "align": PP_ALIGN.CENTER}])
    add_text_box(slide, Inches(0.42), Inches(3.9), Inches(12.5), Inches(0.6),
                 [{"text": "Integrazione One Fiscale su TCF",
                   "size": Pt(20), "color": WHITE, "align": PP_ALIGN.CENTER}])
    add_text_box(slide, Inches(0.42), Inches(4.6), Inches(12.5), Inches(0.5),
                 [{"text": "Wolters Kluwer · NTT DATA – Luglio 2026",
                   "size": Pt(13), "color": BLUE_LIGHT, "align": PP_ALIGN.CENTER}])
    add_text_box(slide, Inches(0.42), Inches(7.1), Inches(12.5), Inches(0.3),
                 [{"text": "© 2026 NTT DATA, Inc.",
                   "size": Pt(9), "color": BLUE_LIGHT, "align": PP_ALIGN.CENTER}])
    # NTT DATA logo bottom-right (white on blue background)
    add_ntt_logo_on_dark(slide, right_x=SLIDE_W - Inches(0.30),
                         bottom_y=Inches(7.20))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def build():
    tmp_src = Path(os.environ["TEMP"]) / "template_probe.pptx"
    if not tmp_src.exists() or tmp_src.stat().st_size != SRC.stat().st_size:
        shutil.copyfile(SRC, tmp_src)
    tmp_out = Path(os.environ["TEMP"]) / "tcf_new.pptx"
    shutil.copyfile(tmp_src, tmp_out)

    prs = Presentation(tmp_out)
    clear_slides(prs)

    build_slide_1_cover(prs)
    build_slide_2_context(prs)
    build_slide_3_objective(prs)
    build_slide_4_architecture(prs)
    build_slide_6_identity(prs)
    build_slide_7_poc(prs)
    build_slide_8_costs(prs)
    build_slide_9_end(prs)

    prs.save(tmp_out)
    shutil.copyfile(tmp_out, OUT)
    print(f"Wrote: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
