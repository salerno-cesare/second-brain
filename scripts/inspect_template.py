"""Inspect the Unicredit RAG Proposal template to understand its structure."""
import os
from pptx import Presentation
from pptx.util import Emu

src = os.path.join(os.environ["TEMP"], "uc_template.pptx")
prs = Presentation(src)

print(f"Slide size: {Emu(prs.slide_width).inches}in x {Emu(prs.slide_height).inches}in")
print(f"Slide count: {len(prs.slides)}")
print(f"Slide layouts count: {len(prs.slide_layouts)}")
print()

print("=== SLIDE LAYOUTS ===")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: {layout.name!r}")
    for ph in layout.placeholders:
        print(f"     ph idx={ph.placeholder_format.idx} type={ph.placeholder_format.type} name={ph.name!r}")
print()

print("=== SLIDES ===")
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} (layout: {slide.slide_layout.name!r}) ---")
    for shape in slide.shapes:
        info = f"  shape={shape.shape_type} name={shape.name!r}"
        if shape.has_text_frame:
            txt = "\n     | ".join(
                (p.text or "") for p in shape.text_frame.paragraphs
            )
            info += f"\n     TEXT:\n     | {txt}"
        if shape.is_placeholder:
            info += f"  (ph idx={shape.placeholder_format.idx} type={shape.placeholder_format.type})"
        print(info)
